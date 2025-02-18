from abc import ABC
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders.base import BaseLoader

from comps import CustomLogger
from rag.module.indexing.loader.ocr import get_rapid_ocr

logger = CustomLogger("pptx_loader")
# 设置图片 OCR 阈值比例
PPTX_OCR_THRESHOLD = (0.3, 0.3)  # (宽度阈值, 高度阈值)
TEXT_LENGTH_THRESHOLD = 500
class CustomizedPPTXLoader(BaseLoader, ABC):
    def __init__(self, file_path):
        self.file_path = file_path
        self.ocr = get_rapid_ocr()  # OCR 模块

    def load(self):
        ppt = Presentation(self.file_path)
        # 幻灯片的宽和高
        slide_width, slide_height = ppt.slide_width/9525, ppt.slide_height/9525
        txts = []
        for slide_num, slide in enumerate(ppt.slides):
            texts = []
            for shape in sorted(
                    slide.shapes, key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left)):

                txt = self.__extract(shape, slide_num, slide_width, slide_height)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))
        # 创建一个 Document 对象，包含聚合后的内容
        doc = Document(page_content="\n".join(txts), metadata={"source": self.file_path})
        return [doc]  # 返回一个包含文档的列表

    def __extract(self, shape, slide_num, slide_width, slide_height):
        # 处理表格内容
        if shape.shape_type == 19:  # Table 类型
            tb = shape.table
            rows = []
            for i in range(1, len(tb.rows)):
                rows.append("; ".join([tb.cell(0, j).text + ": " + tb.cell(i, j).text
                                       for j in range(len(tb.columns)) if tb.cell(i, j)]))
            return "\n".join(rows)

        # 处理文本框内容
        if shape.has_text_frame:
            return shape.text_frame.text

        # 处理 Group 类型内容
        if shape.shape_type == 6:
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                t = self.__extract(p, slide_num, slide_width, slide_height)
                if t:
                    texts.append(t)
            return "\n".join(texts)

        # 处理图片内容并进行 OCR
        if shape.shape_type == 13:  # Picture 类型
            image = shape.image
            image_bytes = image.blob

            # 获取图片和幻灯片尺寸
            display_width, display_height = shape.width / 9525, shape.height / 9525
            width_ratio = display_width / slide_width
            height_ratio = display_height / slide_height
            # # 判断图片大小是否超过阈值
            if width_ratio >= PPTX_OCR_THRESHOLD[0] or height_ratio >= PPTX_OCR_THRESHOLD[1]:
                # 保存图片到本地
                result, _ = self.ocr(image_bytes)
                if result:
                    ocr_result = [line[1] for line in result]
                    logger.info("ocr result:", ocr_result)
                    return "\n".join(ocr_result)
        return None  # 如果不满足阈值条件，则返回 None



# 使用示例
if __name__ == "__main__":
    pptx_file_path = "./data/raw/pptx/example.pptx"
    pptx_loader = CustomizedPPTXLoader(pptx_file_path)  # 替换为你的 .pptx 文件路径
    documents = pptx_loader.load()
    for doc in documents:
        print(doc.page_content)
